import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
from typing import Literal
import os
from openpyxl import Workbook
import pandas as pd
from ..utils import PptConfig
from typing import Optional


def df_to_excel(df: pd.DataFrame, excel_path: str, sheet_name: str):
    if df.shape == (0, 0):
        print('Current state of dataframe is None')
    if not os.path.exists(excel_path):
        workbook = Workbook()
        sheet = workbook.active  # Truy cập vào sheet hiện tại
        sheet.title = sheet_name  # Đặt tên cho sheet
        workbook.save(excel_path)
    with pd.ExcelWriter(excel_path, engine='openpyxl', mode='a', if_sheet_exists='replace') as writer:
        df.to_excel(writer, sheet_name=sheet_name)

chart_type_lookup = {
    'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
    'bar': XL_CHART_TYPE.BAR_CLUSTERED,
    'pie': XL_CHART_TYPE.PIE
}

def create_pptx_chart(
        template_path,
        dataframe,
        type=Literal['column', 'bar', 'pie'], 
        title=None,
        config: Optional[PptConfig] = None
    ):
    
    if not config:
        config = PptConfig()
    
    theme_color = config.theme_color
    
    prs = Presentation(template_path) if os.path.exists(template_path) else Presentation()
    chart_type = chart_type_lookup[type]
    slide_layout = prs.slide_layouts[config.slide_layout]  # Layout trống
    slide = prs.slides.add_slide(slide_layout)

    if type == 'pie':
        chart_data = ChartData()
        chart_data.categories = [str(i) for i in dataframe['category']]
        chart_data.add_series('Series', dataframe['count'])
    else:
        chart_data = CategoryChartData()
        if 'category' in dataframe.columns:
            chart_data.categories = [str(i) for i in dataframe['category']]
            chart_data.add_series('count', dataframe['count'])
        elif 'row_value' in dataframe.columns:
            chart_data.categories = dataframe['row_value']
            for series_name in dataframe.columns[1:]:
                chart_data.add_series(series_name, dataframe[series_name])

    x, y, cx, cy = Inches(config.position[0]), Inches(config.position[1]), Inches(config.position[2]), Inches(config.position[3])
    chart = slide.shapes.add_chart(
        chart_type,  # Sử dụng loại biểu đồ này
        x, y, cx, cy,
        chart_data
    ).chart

    chart.font.name = config.font
    chart.has_legend = config.has_legend
    chart.has_title = config.has_title
    chart.chart_title.text_frame.text = title
    chart.legend.position = config.legend_position
    chart.legend.font.size = Pt(config.legend_font_size)
    try:
        chart.category_axis.has_major_gridlines = config.category_axis_has_major_gridlines
        chart.category_axis.has_minor_gridlines = config.category_axis_has_minor_gridlines
        chart.category_axis.has_title = config.category_axis_has_title
        chart.category_axis.visible = config.category_axis_visible
        chart.category_axis.tick_labels.font.size = Pt(config.category_axis_tick_labels_font_size)
    except:
        pass
    try:
        chart.value_axis.has_major_gridlines = config.value_axis_has_major_gridlines
        chart.value_axis.has_minor_gridlines = config.value_axis_has_minor_gridlines
        chart.value_axis.visible = config.value_axis_visible
    except:
        pass

    if type != 'pie':
        for index, i in enumerate(chart.series):
            i.data_labels.font.size = Pt(config.data_labels_font_size)
            i.data_labels.font.name = config.data_labels_font
            i.data_labels.number_format = config.data_labels_number_format
            i.data_labels.number_format_is_linked = config.data_labels_number_format_is_linked
            i.data_labels.position = config.data_labels_position
            i.data_labels.show_category_name = config.data_labels_show_category_name
            i.data_labels.show_legend_key = config.data_labels_show_legend_key
            i.data_labels.show_percentage = config.data_labels_show_percentage
            i.data_labels.show_series_name = config.data_labels_show_series_name
            i.data_labels.show_value = config.data_labels_show_value
            i.format.fill.solid()
            try:
                i.format.fill.fore_color.theme_color = theme_color[index]
            except:
                i.format.fill.fore_color.theme_color = theme_color[0]
    else:
        for serie in chart.series:
            for idx, point in enumerate(serie.points):
                point.format.fill.solid()
                try:
                    point.format.fill.fore_color.theme_color = theme_color[idx]
                except:
                    point.format.fill.fore_color.theme_color = theme_color[0]

    prs.save(template_path)

